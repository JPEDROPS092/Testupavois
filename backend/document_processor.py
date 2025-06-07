#!/usr/bin/env python3
"""
Document processor module for Testupavois TTS application.
Handles conversion of various document formats to formatted text.
"""

import os
import re
import io
from typing import Dict, Any, Optional, List, Tuple

# Document processing libraries
import docx
import PyPDF2
from pptx import Presentation

class DocumentProcessor:
    """
    Processes various document formats and extracts formatted text for TTS.
    Supported formats: PDF, DOCX, PPTX, TXT
    """
    
    @staticmethod
    def process_document(file_content: bytes, filename: str) -> str:
        """
        Process a document file and extract formatted text.
        
        Args:
            file_content: The binary content of the file
            filename: Original filename with extension
            
        Returns:
            Formatted text extracted from the document
        
        Raises:
            ValueError: If file format is not supported or processing fails
        """
        file_ext = os.path.splitext(filename)[1].lower()
        
        if file_ext == '.pdf':
            return DocumentProcessor._process_pdf(file_content)
        elif file_ext == '.docx':
            return DocumentProcessor._process_docx(file_content)
        elif file_ext == '.pptx':
            return DocumentProcessor._process_pptx(file_content)
        elif file_ext == '.txt':
            return DocumentProcessor._process_txt(file_content)
        else:
            raise ValueError(f"Unsupported file format: {file_ext}")
    
    @staticmethod
    def _process_pdf(file_content: bytes) -> str:
        """Extract and format text from PDF files."""
        try:
            pdf_file = io.BytesIO(file_content)
            reader = PyPDF2.PdfReader(pdf_file)
            
            text_parts = []
            for i, page in enumerate(reader.pages):
                page_text = page.extract_text()
                if page_text.strip():
                    text_parts.append(f"Página {i+1}.\n{page_text.strip()}\n")
            
            return "\n".join(text_parts)
        except Exception as e:
            raise ValueError(f"Failed to process PDF: {str(e)}")
    
    @staticmethod
    def _process_docx(file_content: bytes) -> str:
        """Extract and format text from DOCX files."""
        try:
            docx_file = io.BytesIO(file_content)
            doc = docx.Document(docx_file)
            
            text_parts = []
            
            # Process paragraphs with formatting awareness
            for para in doc.paragraphs:
                if para.text.strip():
                    # Check if paragraph is a heading
                    if para.style.name.startswith('Heading'):
                        heading_level = int(para.style.name.replace('Heading', '')) if para.style.name != 'Heading' else 1
                        text_parts.append(f"{'#' * heading_level} {para.text.strip()}")
                    else:
                        text_parts.append(para.text.strip())
            
            # Process tables
            for table in doc.tables:
                table_text = []
                for row in table.rows:
                    row_text = [cell.text.strip() for cell in row.cells]
                    table_text.append(" | ".join(row_text))
                
                if table_text:
                    text_parts.append("Tabela:\n" + "\n".join(table_text))
            
            return "\n\n".join(text_parts)
        except Exception as e:
            raise ValueError(f"Failed to process DOCX: {str(e)}")
    
    @staticmethod
    def _process_pptx(file_content: bytes) -> str:
        """Extract and format text from PPTX files."""
        try:
            pptx_file = io.BytesIO(file_content)
            presentation = Presentation(pptx_file)
            
            text_parts = []
            
            for i, slide in enumerate(presentation.slides):
                slide_text = []
                slide_text.append(f"Slide {i+1}.")
                
                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())
                
                if len(slide_text) > 1:  # If we have more than just the slide number
                    text_parts.append("\n".join(slide_text))
            
            return "\n\n".join(text_parts)
        except Exception as e:
            raise ValueError(f"Failed to process PPTX: {str(e)}")
    
    @staticmethod
    def _process_txt(file_content: bytes) -> str:
        """Process plain text files."""
        try:
            # Try different encodings
            encodings = ['utf-8', 'latin-1', 'cp1252']
            text = None
            
            for encoding in encodings:
                try:
                    text = file_content.decode(encoding)
                    break
                except UnicodeDecodeError:
                    continue
            
            if text is None:
                raise ValueError("Could not decode text file with any of the attempted encodings")
                
            # Clean up the text
            text = re.sub(r'\r\n', '\n', text)  # Normalize line endings
            text = re.sub(r'\n{3,}', '\n\n', text)  # Remove excessive line breaks
            
            return text
        except Exception as e:
            raise ValueError(f"Failed to process TXT: {str(e)}")
    
    @staticmethod
    def format_text_for_tts(text: str, max_length: int = 5000) -> List[str]:
        """
        Format and split text into chunks suitable for TTS processing.
        
        Args:
            text: The input text to format
            max_length: Maximum length of each chunk
            
        Returns:
            List of text chunks ready for TTS
        """
        # Clean up text
        text = re.sub(r'\s+', ' ', text)  # Replace multiple spaces with single space
        
        # Split by sentences to avoid cutting in the middle of a sentence
        sentence_pattern = r'(?<=[.!?])\s+'
        sentences = re.split(sentence_pattern, text)
        
        chunks = []
        current_chunk = ""
        
        for sentence in sentences:
            # If adding this sentence would exceed max_length, start a new chunk
            if len(current_chunk) + len(sentence) > max_length:
                if current_chunk:
                    chunks.append(current_chunk.strip())
                current_chunk = sentence + " "
            else:
                current_chunk += sentence + " "
        
        # Add the last chunk if not empty
        if current_chunk:
            chunks.append(current_chunk.strip())
        
        return chunks
